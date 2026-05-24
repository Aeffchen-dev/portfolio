"""
Frontify "Rebranding Redefined" — TRUE Slide Master
────────────────────────────────────────────────────
Creates a .pptx with a fully configured Slide Master:
  • Theme: Frontify brand colors (6 accents) + Montserrat / Inter fonts
  • Master slide: dark background, left accent bar, wordmark footer
  • 11 proper slide layouts with title / content placeholders

Import into Google Slides via:
  File → Import theme  (applies colors + fonts to all future slides)
  — OR —
  File → Open (uploads full deck; layouts available via "Apply layout")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import lxml.etree as etree
import copy
import io

# ─────────────────────────────────────────────────────────────
# NAMESPACE CONSTANTS
# ─────────────────────────────────────────────────────────────
NS_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P  = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"

def a(tag):  return f"{{{NS_A}}}{tag}"
def p(tag):  return f"{{{NS_P}}}{tag}"

# ─────────────────────────────────────────────────────────────
# BRAND TOKENS
# ─────────────────────────────────────────────────────────────
HEX = dict(
    dark_bg    = "100E1D",
    violet     = "5834E8",
    violet_dk  = "3A20B0",
    violet_lt  = "C4B4FF",
    cream      = "F2EDE0",
    cream_mid  = "E3DBC8",
    white      = "FFFFFF",
    dark_text  = "1A162A",
    mid_gray   = "8884A0",
    light_gray = "C8C4D8",
    # chapter accents
    ch01 = "5834E8",  # Typography  – violet
    ch02 = "FF6038",  # Color       – orange
    ch03 = "00D4B0",  # Sound       – teal
    ch04 = "3880FF",  # Motion      – blue
    ch05 = "30D158",  # Flexibility – green
    ch06 = "FF375F",  # Culture     – crimson
)

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# ─────────────────────────────────────────────────────────────
# SLIDE DIMENSIONS
# ─────────────────────────────────────────────────────────────
W = Inches(13.333)   # 16 : 9 widescreen
H = Inches(7.5)
M = Inches(0.9)      # standard margin

# ─────────────────────────────────────────────────────────────
# THEME XML PATCH
# ─────────────────────────────────────────────────────────────
THEME_COLORS = {
    "dk1":     ("srgbClr", HEX["dark_bg"]),   # Dark 1  → text default
    "lt1":     ("srgbClr", HEX["cream"]),     # Light 1 → bg default
    "dk2":     ("srgbClr", HEX["dark_text"]),
    "lt2":     ("srgbClr", HEX["white"]),
    "accent1": ("srgbClr", HEX["violet"]),    # Accent 1
    "accent2": ("srgbClr", HEX["ch02"]),      # Accent 2
    "accent3": ("srgbClr", HEX["ch03"]),      # Accent 3
    "accent4": ("srgbClr", HEX["ch04"]),      # Accent 4
    "accent5": ("srgbClr", HEX["ch05"]),      # Accent 5
    "accent6": ("srgbClr", HEX["ch06"]),      # Accent 6
    "hlink":   ("srgbClr", HEX["violet"]),
    "folHlink":("srgbClr", HEX["violet_dk"]),
}

def patch_theme(theme_blob: bytes) -> bytes:
    """Rewrite theme colors and fonts in the theme XML blob."""
    root = etree.fromstring(theme_blob)
    root.set("name", "Frontify")

    # ── colors ──
    clr_scheme = root.find(f".//{a('clrScheme')}")
    if clr_scheme is not None:
        clr_scheme.set("name", "Frontify")
        for role, (elem_name, val) in THEME_COLORS.items():
            node = clr_scheme.find(a(role))
            if node is None:
                node = etree.SubElement(clr_scheme, a(role))
            for child in list(node):
                node.remove(child)
            child = etree.SubElement(node, a(elem_name))
            child.set("val", val)

    # ── fonts ──
    font_scheme = root.find(f".//{a('fontScheme')}")
    if font_scheme is not None:
        font_scheme.set("name", "Frontify")
        for kind, typeface in (("majorFont", "Montserrat"), ("minorFont", "Inter")):
            fnode = font_scheme.find(a(kind))
            if fnode is not None:
                latin = fnode.find(a("latin"))
                if latin is None:
                    latin = etree.SubElement(fnode, a("latin"))
                latin.set("typeface", typeface)

    return etree.tostring(root, xml_declaration=True,
                          encoding="UTF-8", standalone=True)


# ─────────────────────────────────────────────────────────────
# LOW-LEVEL SHAPE XML BUILDERS
# ─────────────────────────────────────────────────────────────

def _sp_id(n: int) -> int:
    return 1000 + n   # avoid collisions with existing placeholder IDs


def solid_rect_xml(shape_id, name, x, y, cx, cy, hex_fill):
    """Return <p:sp> XML for a solid-fill rectangle (no text)."""
    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{hex_fill}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>'''


def text_rect_xml(shape_id, name, x, y, cx, cy,
                  text, font, size_pt, bold, hex_color,
                  align="l", hex_fill=None, italic=False):
    """Return <p:sp> XML for a rectangle with styled text."""
    bold_attr  = "1" if bold   else "0"
    italic_attr= "1" if italic else "0"
    fill_xml = (
        f"<a:solidFill><a:srgbClr val=\"{hex_fill}\"/></a:solidFill>"
        if hex_fill else "<a:noFill/>"
    )
    # Escape text
    safe = (text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                .replace('"', "&quot;").replace("'","&apos;"))
    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"/>
    <p:cNvSpPr txBox="1"><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    {fill_xml}
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="0" rIns="0" tIns="0" bIns="0"/>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="{align}"/>
      <a:r>
        <a:rPr lang="en-US" sz="{int(size_pt*100)}" b="{bold_attr}"
               i="{italic_attr}" dirty="0">
          <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
          <a:latin typeface="{font}"/>
        </a:rPr>
        <a:t>{safe}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def placeholder_xml(ph_type, ph_idx, shape_id, name,
                    x, y, cx, cy,
                    font_major, size_pt, bold, hex_color, align="l",
                    hex_bg=None):
    """Return a <p:sp> placeholder element."""
    bold_attr = "1" if bold else "0"
    fill_xml = (
        f"<a:solidFill><a:srgbClr val=\"{hex_bg}\"/></a:solidFill>"
        if hex_bg else "<a:noFill/>"
    )
    return f'''<p:sp xmlns:p="{NS_P}" xmlns:a="{NS_A}" xmlns:r="{NS_R}">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/>
      <p:ph type="{ph_type}" idx="{ph_idx}"/>
    </p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{int(x)}" y="{int(y)}"/><a:ext cx="{int(cx)}" cy="{int(cy)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    {fill_xml}
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle>
      <a:lvl1pPr algn="{align}">
        <a:defRPr sz="{int(size_pt*100)}" b="{bold_attr}" dirty="0">
          <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
          <a:latin typeface="+mj-lt" pitchFamily="0" charset="0"/>
        </a:defRPr>
      </a:lvl1pPr>
    </a:lstStyle>
    <a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>
  </p:txBody>
</p:sp>'''


def inject_spTree(spTree, xml_str):
    """Parse xml_str and append its root into spTree."""
    elem = etree.fromstring(xml_str)
    spTree.append(elem)


# ─────────────────────────────────────────────────────────────
# MASTER BACKGROUND HELPER
# ─────────────────────────────────────────────────────────────

def set_master_background(master_elem, hex_color):
    """Set <p:cSld><p:bg> on a slide master element."""
    cSld = master_elem.find(p("cSld"))
    if cSld is None:
        return
    # Remove existing bg
    for bg in cSld.findall(p("bg")):
        cSld.remove(bg)
    bg_xml = f'''<p:bg xmlns:p="{NS_P}" xmlns:a="{NS_A}">
  <p:bgPr>
    <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
    <a:effectLst/>
  </p:bgPr>
</p:bg>'''
    cSld.insert(0, etree.fromstring(bg_xml))


def set_layout_background(layout_elem, hex_color):
    """Set <p:cSld><p:bg> on a slide layout element."""
    set_master_background(layout_elem, hex_color)   # same structure


# ─────────────────────────────────────────────────────────────
# CLEAR ALL SHAPES FROM A LAYOUT
# ─────────────────────────────────────────────────────────────

def clear_spTree(elem):
    """Remove all sp / grpSp / graphicFrame from spTree, keep nvGrpSpPr/grpSpPr."""
    cSld = elem.find(p("cSld"))
    if cSld is None:
        return
    spTree = cSld.find(p("spTree"))
    if spTree is None:
        return
    keep_tags = {p("nvGrpSpPr"), p("grpSpPr")}
    for child in list(spTree):
        if child.tag not in keep_tags:
            spTree.remove(child)
    return spTree


# ─────────────────────────────────────────────────────────────
# HELPER: inches / points → EMU
# ─────────────────────────────────────────────────────────────
def i(n):   return int(Inches(n))
def pt(n):  return int(Pt(n))

# ─────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ─────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

master = prs.slide_masters[0]

# ── 1. PATCH THEME ────────────────────────────────────────────
for rid, rel in master.part.rels.items():
    if "theme" in rel.reltype:
        tp = rel.target_part
        tp._blob = patch_theme(tp.blob)
        break

# ── 2. MASTER SLIDE: dark background + brand shapes ───────────
master_elem = master.element
set_master_background(master_elem, HEX["dark_bg"])

cSld_m   = master_elem.find(p("cSld"))
spTree_m = cSld_m.find(p("spTree"))

# Remove decorative shapes only (keep placeholders / type elements)
for child in list(spTree_m):
    tag = child.tag
    if tag == p("sp"):
        ph = child.find(f".//{p('ph')}")
        if ph is None:                  # not a placeholder → remove
            spTree_m.remove(child)

# Left accent bar (appears on every slide via master inheritance)
inject_spTree(spTree_m, solid_rect_xml(
    _sp_id(1), "AccentBar",
    0, 0, i(0.22), i(7.5), HEX["violet"]))

# Wordmark bottom-right
inject_spTree(spTree_m, text_rect_xml(
    _sp_id(2), "Wordmark",
    i(11.3), i(7.08), i(1.8), i(0.38),
    "frontify", "Montserrat", 13, True, HEX["violet"]))

# ── 3. CONFIGURE SLIDE LAYOUTS ────────────────────────────────
layouts = master.slide_layouts
# Index from inspection:
# 0  Title Slide          → Cover
# 1  Title and Content    → Standard content
# 2  Section Header       → Chapter divider
# 3  Two Content          → Two-column
# 4  Comparison           → Comparison
# 5  Title Only           → Title only
# 6  Blank                → Keep blank
# 7  Content with Caption → Quote / callout
# 8  Picture with Caption → Full-bleed image
# 9  Title and Vertical Text → Stats
# 10 Vertical Title and Text → Summary

# ─────────────────────────────────────────────────────────────
# LAYOUT 0: COVER / TITLE (dark)
# ─────────────────────────────────────────────────────────────
ly = layouts[0]
ly.name = "01 – Cover"
set_layout_background(ly.element, HEX["dark_bg"])
spTree = clear_spTree(ly.element)

# Title placeholder
inject_spTree(spTree, placeholder_xml(
    "title", 0, 10, "Title",
    i(0.9), i(1.1), i(9.5), i(3.2),
    "Montserrat", 72, True, HEX["white"]))

# Subtitle placeholder
inject_spTree(spTree, placeholder_xml(
    "subTitle", 1, 11, "Subtitle",
    i(0.9), i(4.55), i(9.0), i(0.7),
    "Inter", 22, False, HEX["light_gray"]))

# Violet rule
inject_spTree(spTree, solid_rect_xml(
    _sp_id(10), "Rule",
    i(0.9), i(4.4), i(2.4), pt(3), HEX["violet"]))

# Date/author line
inject_spTree(spTree, text_rect_xml(
    _sp_id(11), "DateLine",
    i(0.9), i(5.55), i(6), i(0.5),
    "A Frontify Guide  ·  2024", "Inter", 14, False, HEX["mid_gray"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 1: TITLE AND CONTENT (cream)
# ─────────────────────────────────────────────────────────────
ly = layouts[1]
ly.name = "02 – Title and Content"
set_layout_background(ly.element, HEX["cream"])
spTree = clear_spTree(ly.element)

inject_spTree(spTree, placeholder_xml(
    "title", 0, 20, "Title",
    i(0.9), i(0.55), i(11.5), i(0.9),
    "Montserrat", 36, True, HEX["dark_text"]))

# Divider rule
inject_spTree(spTree, solid_rect_xml(
    _sp_id(21), "Rule",
    i(0.9), i(1.55), i(11.5), pt(1), HEX["cream_mid"]))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 22, "Content",
    i(0.9), i(1.75), i(11.5), i(5.0),
    "Inter", 18, False, HEX["dark_text"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 2: SECTION HEADER / CHAPTER DIVIDER (dark + violet block)
# ─────────────────────────────────────────────────────────────
ly = layouts[2]
ly.name = "03 – Section Header"
set_layout_background(ly.element, HEX["dark_bg"])
spTree = clear_spTree(ly.element)

# Colour block (left half) — uses Accent 1 (violet) via theme
inject_spTree(spTree, solid_rect_xml(
    _sp_id(30), "ColorBlock",
    0, 0, i(5.8), i(7.5), HEX["violet"]))

# Chapter number (giant ghost)
inject_spTree(spTree, text_rect_xml(
    _sp_id(31), "ChapterNum",
    i(0.1), i(-0.8), i(5.6), i(6.5),
    "01", "Montserrat", 260, True, "7B5CF0"))

# Right side label
inject_spTree(spTree, text_rect_xml(
    _sp_id(32), "ChapterLabel",
    i(6.3), i(1.4), i(6.5), i(0.5),
    "CHAPTER 01", "Inter", 13, False, HEX["violet_lt"]))

inject_spTree(spTree, placeholder_xml(
    "title", 0, 33, "ChapterTitle",
    i(6.3), i(1.95), i(6.5), i(2.2),
    "Montserrat", 52, True, HEX["white"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(34), "AccentRule",
    i(6.3), i(4.25), i(2.0), pt(2), HEX["violet_lt"]))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 35, "ChapterDesc",
    i(6.3), i(4.45), i(6.5), i(1.5),
    "Inter", 18, False, HEX["light_gray"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 3: TWO CONTENT (dark)
# ─────────────────────────────────────────────────────────────
ly = layouts[3]
ly.name = "04 – Two Column"
set_layout_background(ly.element, HEX["dark_bg"])
spTree = clear_spTree(ly.element)

inject_spTree(spTree, placeholder_xml(
    "title", 0, 40, "Title",
    i(0.9), i(0.55), i(11.5), i(0.85),
    "Montserrat", 34, True, HEX["white"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(41), "Rule",
    i(0.9), i(1.5), i(11.5), pt(1), "38345A"))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 42, "LeftContent",
    i(0.9), i(1.7), i(5.5), i(5.3),
    "Inter", 17, False, HEX["light_gray"]))

inject_spTree(spTree, placeholder_xml(
    "body", 2, 43, "RightContent",
    i(7.1), i(1.7), i(5.5), i(5.3),
    "Inter", 17, False, HEX["light_gray"]))

# Vertical divider
inject_spTree(spTree, solid_rect_xml(
    _sp_id(44), "ColDivider",
    i(6.72), i(1.7), pt(1), i(5.3), "38345A"))

# ─────────────────────────────────────────────────────────────
# LAYOUT 4: CONTENT + IMAGE RIGHT (cream)
# ─────────────────────────────────────────────────────────────
ly = layouts[4]
ly.name = "05 – Content + Image (right)"
set_layout_background(ly.element, HEX["cream"])
spTree = clear_spTree(ly.element)

inject_spTree(spTree, placeholder_xml(
    "title", 0, 50, "Title",
    i(0.9), i(1.1), i(5.8), i(1.8),
    "Montserrat", 38, True, HEX["dark_text"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(51), "Rule",
    i(0.9), i(3.05), i(1.4), pt(2), HEX["violet"]))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 52, "BodyText",
    i(0.9), i(3.2), i(5.8), i(3.1),
    "Inter", 17, False, HEX["dark_text"]))

# Image placeholder (picture)
inject_spTree(spTree, solid_rect_xml(
    _sp_id(53), "ImageArea",
    i(7.1), 0, i(6.233), i(7.5), "D8D2C4"))

inject_spTree(spTree, text_rect_xml(
    _sp_id(54), "ImageLabel",
    i(8.5), i(3.45), i(3.5), i(0.6),
    "[IMAGE]", "Inter", 12, False, HEX["mid_gray"],
    align="ctr"))

# ─────────────────────────────────────────────────────────────
# LAYOUT 5: TITLE ONLY (dark)
# ─────────────────────────────────────────────────────────────
ly = layouts[5]
ly.name = "06 – Title Only"
set_layout_background(ly.element, HEX["dark_bg"])
spTree = clear_spTree(ly.element)

inject_spTree(spTree, placeholder_xml(
    "title", 0, 60, "Title",
    i(0.9), i(0.55), i(11.5), i(0.85),
    "Montserrat", 36, True, HEX["white"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(61), "Rule",
    i(0.9), i(1.55), i(11.5), pt(1), "38345A"))

# ─────────────────────────────────────────────────────────────
# LAYOUT 6: BLANK (keep clean)
# ─────────────────────────────────────────────────────────────
ly = layouts[6]
ly.name = "07 – Blank (dark)"
set_layout_background(ly.element, HEX["dark_bg"])
clear_spTree(ly.element)

# ─────────────────────────────────────────────────────────────
# LAYOUT 7: PULL QUOTE (dark, centered)
# ─────────────────────────────────────────────────────────────
ly = layouts[7]
ly.name = "08 – Pull Quote"
set_layout_background(ly.element, HEX["dark_bg"])
spTree = clear_spTree(ly.element)

# Decorative quote mark
inject_spTree(spTree, text_rect_xml(
    _sp_id(70), "QuoteMark",
    i(0.5), i(-0.5), i(3), i(2.5),
    "“", "Montserrat", 200, True, HEX["violet"]))

inject_spTree(spTree, placeholder_xml(
    "body", 0, 71, "QuoteText",
    i(0.9), i(1.8), i(11.5), i(2.9),
    "Montserrat", 44, True, HEX["white"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(72), "AttrRule",
    i(0.9), i(4.9), i(0.8), pt(2), HEX["violet"]))

inject_spTree(spTree, placeholder_xml(
    "title", 1, 73, "Attribution",
    i(0.9), i(5.1), i(9), i(0.6),
    "Inter", 15, False, HEX["mid_gray"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 8: FULL-BLEED IMAGE WITH OVERLAY
# ─────────────────────────────────────────────────────────────
ly = layouts[8]
ly.name = "09 – Full-Bleed Image"
set_layout_background(ly.element, "18142A")
spTree = clear_spTree(ly.element)

# Full-slide image placeholder zone
inject_spTree(spTree, solid_rect_xml(
    _sp_id(80), "ImageFill",
    0, 0, i(13.333), i(7.5), "221C34"))

inject_spTree(spTree, text_rect_xml(
    _sp_id(81), "ImageLabel",
    i(5.5), i(3.35), i(2.5), i(0.55),
    "[IMAGE]", "Inter", 12, False, HEX["mid_gray"], align="ctr"))

# Gradient overlay (solid dark at bottom)
inject_spTree(spTree, solid_rect_xml(
    _sp_id(82), "Overlay",
    0, i(4.0), i(13.333), i(3.5), "08060E"))

inject_spTree(spTree, placeholder_xml(
    "title", 0, 83, "Caption",
    i(0.9), i(4.3), i(11.5), i(1.7),
    "Montserrat", 40, True, HEX["white"]))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 84, "Credit",
    i(0.9), i(6.3), i(9), i(0.55),
    "Inter", 13, False, HEX["mid_gray"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 9: KEY INSIGHT / CALLOUT (violet)
# ─────────────────────────────────────────────────────────────
ly = layouts[9]
ly.name = "10 – Key Insight"
set_layout_background(ly.element, HEX["violet"])
spTree = clear_spTree(ly.element)

# Corner block decoration
inject_spTree(spTree, solid_rect_xml(
    _sp_id(90), "Corner",
    i(10.5), i(-0.6), i(3.0), i(3.0), HEX["violet_dk"]))

inject_spTree(spTree, text_rect_xml(
    _sp_id(91), "Eyebrow",
    i(0.9), i(1.5), i(10), i(0.5),
    "KEY INSIGHT", "Inter", 13, False, HEX["violet_lt"]))

inject_spTree(spTree, placeholder_xml(
    "title", 0, 92, "InsightText",
    i(0.9), i(2.1), i(10.5), i(2.5),
    "Montserrat", 52, True, HEX["white"]))

inject_spTree(spTree, solid_rect_xml(
    _sp_id(93), "Rule",
    i(0.9), i(4.75), i(1.8), pt(2), HEX["violet_lt"]))

inject_spTree(spTree, placeholder_xml(
    "body", 1, 94, "Context",
    i(0.9), i(5.0), i(9), i(0.6),
    "Inter", 16, False, HEX["violet_lt"]))

# ─────────────────────────────────────────────────────────────
# LAYOUT 10: STATISTICS (cream)
# ─────────────────────────────────────────────────────────────
ly = layouts[10]
ly.name = "11 – Statistics"
set_layout_background(ly.element, HEX["cream"])
spTree = clear_spTree(ly.element)

inject_spTree(spTree, text_rect_xml(
    _sp_id(100), "Eyebrow",
    i(0.9), i(0.5), i(10), i(0.5),
    "BY THE NUMBERS", "Inter", 13, False, HEX["mid_gray"]))

# Three stat columns
stat_x = [0.9, 5.05, 9.2]
accent_hexes = [HEX["ch02"], HEX["ch03"], HEX["ch04"]]
labels = ["Stat 1", "Stat 2", "Stat 3"]
for j, (x, acc, lbl) in enumerate(zip(stat_x, accent_hexes, labels)):
    inject_spTree(spTree, solid_rect_xml(
        _sp_id(101 + j*4), f"AccentBar{j}",
        i(x), i(1.25), i(3.7), pt(4), acc))
    inject_spTree(spTree, text_rect_xml(
        _sp_id(102 + j*4), f"BigNum{j}",
        i(x), i(1.4), i(3.7), i(2.0),
        "00%", "Montserrat", 80, True, HEX["dark_text"]))
    inject_spTree(spTree, solid_rect_xml(
        _sp_id(103 + j*4), f"DivLine{j}",
        i(x), i(3.5), i(3.7), pt(1), HEX["cream_mid"]))
    inject_spTree(spTree, text_rect_xml(
        _sp_id(104 + j*4), f"StatDesc{j}",
        i(x), i(3.65), i(3.7), i(1.8),
        "Description of the statistic.", "Inter", 16, False, HEX["dark_text"]))

# ─────────────────────────────────────────────────────────────
# EXAMPLE SLIDES (one per layout)
# ─────────────────────────────────────────────────────────────
blank_layout = layouts[6]  # use the blank layout for example slides

EXAMPLE_DATA = [
    (0, "REBRANDING\nREDEFINED",   "Six elements of modern brand craft"),
    (1, "Typography as Voice",     "How typefaces shape perception and brand personality "
                                   "across every touchpoint."),
    (2, "CHAPTER 01",              "Typography\n\nHow letterforms shape brand voice "
                                   "and personality."),
    (3, "Headline Spanning Width", "Left column content.\n\n• Key point one\n• Key point two"),
    (4, "Section Headline",        "Body copy for the content and image layout."),
    (5, "Standalone Headline",     ""),
    (6, "",                        ""),  # blank
    (7, "Branding is a living system,\nnot a static artifact.",
       "Expert Name  ·  Role, Organisation"),
    (8, "Image Caption or Key Statement Here",
       "Image credit or supporting line"),
    (9, "A successful rebrand is not a\nmoment — it is a system.",
       "Chapter synthesis · Applies across all six elements"),
    (10, "78%", "of brand teams plan a full rebrand within 3 years"),
]

for layout_idx, title_text, body_text in EXAMPLE_DATA:
    layout = layouts[layout_idx]
    slide  = prs.slides.add_slide(layout)

    # Fill placeholders where they exist
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0 and title_text:
            ph.text = title_text
        elif idx == 1 and body_text:
            ph.text = body_text
        elif idx in (2, 3) and body_text:
            ph.text = body_text

# ─────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────
OUT = "/home/user/portfolio/frontify-slide-master.pptx"
prs.save(OUT)
print(f"✓  Saved: {OUT}")
print(f"   Layouts in master: {len(prs.slide_masters[0].slide_layouts)}")
print(f"   Example slides:    {len(prs.slides)}")
print()
print("LAYOUT INDEX")
print("─" * 48)
for i_l, ly in enumerate(prs.slide_masters[0].slide_layouts):
    print(f"  {i_l:02d}  {ly.name}")
