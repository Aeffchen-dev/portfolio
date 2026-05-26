from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
BG     = RGBColor(0x0E, 0x0E, 0x11)
CARD   = RGBColor(0x1C, 0x1C, 0x22)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
VIOLET = RGBColor(0x9E, 0x66, 0xFA)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
YELLOW = RGBColor(0xFA, 0xBF, 0x29)
CORAL  = RGBColor(0xFA, 0x66, 0x5C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x99, 0x99, 0xA4)
LIGHT  = RGBColor(0xD1, 0xD1, 0xD9)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── helpers ─────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def rgb(c: RGBColor):
    return c

def bg_rect(slide, color=BG):
    """Full-slide background rectangle."""
    shp = slide.shapes.add_shape(1, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect(slide, x, y, w, h, fill, radius_pt=0, alpha=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if radius_pt:
        shp.adjustments[0] = radius_pt / 100
    return shp

def bar(slide, x, y, h, color):
    """Thin vertical accent bar."""
    rect(slide, x, y, Pt(5), h, color)

def txt(slide, text, x, y, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def card_rect(slide, x, y, w, h, fill=CARD, radius_pt=8):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def top_bar(slide, x, y, w, color, thickness_pt=4):
    rect(slide, x, y, w, Pt(thickness_pt), color)

def eyebrow(slide, label, x, y, color=VIOLET):
    txt(slide, label, x, y, Inches(5), Pt(28), 16,
        bold=True, color=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 · TITLE
# ═══════════════════════════════════════════════════════════════════════════
s1 = add_slide()
bg_rect(s1)
bar(s1, 0, 0, H, PURPLE)

eyebrow(s1, "Jana Gramlich  ·  2026", Inches(0.7), Inches(0.55), VIOLET)

txt(s1, "Agentic Design Systems",
    Inches(0.7), Inches(1.3), Inches(8.5), Inches(2.8),
    72, bold=True, color=WHITE)

txt(s1,
    "Design infrastructure that thinks, acts,\nand continuously improves itself.",
    Inches(0.7), Inches(4.1), Inches(7), Inches(1.2),
    22, color=MUTED)

# accent rule
rect(s1, Inches(0.7), Inches(5.5), Inches(0.8), Pt(4), PURPLE)

# three theme chips
chips = [("AI Agents", PURPLE), ("Design Systems", VIOLET), ("Automation", GREEN)]
cx = Inches(0.7)
for label, color in chips:
    w = Inches(1.5)
    card_rect(s1, cx, Inches(5.8), w, Pt(30), fill=color)
    txt(s1, label, cx + Pt(8), Inches(5.82), w, Pt(28), 13, bold=True, color=WHITE)
    cx += w + Inches(0.15)

# right orb (three stacked circles via ellipse workaround → use rounded rect)
for r, op_color in [(Inches(1.8), RGBColor(0x2A,0x12,0x5E)),
                    (Inches(1.2), RGBColor(0x4A,0x22,0xAA)),
                    (Inches(0.6), PURPLE)]:
    cx2 = Inches(11.4) - r/2; cy2 = Inches(3.8) - r/2
    shp = slide_ellipse(s1, cx2, cy2, r, r, op_color) if False else None
    shp2 = s1.shapes.add_shape(9, cx2, cy2, r, r)   # freeform → use oval (9=oval)
    shp2.fill.solid(); shp2.fill.fore_color.rgb = op_color
    shp2.line.fill.background()

txt(s1, "⚡", Inches(11.1), Inches(3.5), Inches(0.6), Pt(50), 36, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 · WHAT IS AGENTIC DESIGN?
# ═══════════════════════════════════════════════════════════════════════════
s2 = add_slide()
bg_rect(s2)
bar(s2, 0, 0, H, GREEN)

eyebrow(s2, "What is Agentic Design?", Inches(0.7), Inches(0.55), GREEN)

txt(s2, "A design system that doesn't\nwait to be told what to do.",
    Inches(0.7), Inches(1.2), Inches(9), Inches(2.2),
    52, bold=True, color=WHITE)

# Two contrast cards
# Left — Traditional
lx, cy2, cw, ch = Inches(0.7), Inches(3.8), Inches(5.8), Inches(3.2)
card_rect(s2, lx, cy2, cw, ch, fill=RGBColor(0x26,0x0D,0x0D))
top_bar(s2, lx, cy2, cw, CORAL, 5)
txt(s2, "Traditional Design System", lx+Pt(14), cy2+Pt(16), cw, Pt(28), 15, bold=True, color=CORAL)

old_items = [
    "Published quarterly — drifts immediately",
    "Docs written once, outdated by sprint 2",
    "Humans catch (or miss) every inconsistency",
    "Static components, manual updates",
]
for i, item in enumerate(old_items):
    txt(s2, "✕  " + item,
        lx+Pt(14), cy2+Pt(52)+Pt(i*38), cw-Pt(28), Pt(36),
        13, color=MUTED)

# Right — Agentic
rx = Inches(6.8)
card_rect(s2, rx, cy2, cw, ch, fill=RGBColor(0x0A,0x26,0x18))
top_bar(s2, rx, cy2, cw, GREEN, 5)
txt(s2, "Agentic Design System", rx+Pt(14), cy2+Pt(16), cw, Pt(28), 15, bold=True, color=GREEN)

new_items = [
    "Continuously monitored, self-correcting",
    "Documentation generated from live usage",
    "Agents catch inconsistencies at commit time",
    "Dynamic variants generated on demand",
]
for i, item in enumerate(new_items):
    txt(s2, "✓  " + item,
        rx+Pt(14), cy2+Pt(52)+Pt(i*38), cw-Pt(28), Pt(36),
        13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 · AUTOMATION IDEAS
# ═══════════════════════════════════════════════════════════════════════════
s3 = add_slide()
bg_rect(s3)
bar(s3, 0, 0, H, VIOLET)

eyebrow(s3, "Automation Ideas", Inches(0.7), Inches(0.55), VIOLET)
txt(s3, "What agents can do for your design system",
    Inches(0.7), Inches(1.1), Inches(10), Inches(1.0),
    40, bold=True, color=WHITE)

ideas = [
    ("🔍", "Drift Detection",       "Nightly diff between Figma specs and component code. Auto-raises a PR when tokens or spacing drift.", PURPLE),
    ("♿", "A11y Guardian",          "Scans every new component for contrast, touch target size, and missing ARIA roles before merge.", GREEN),
    ("🎨", "Token Propagation",      "One semantic token update cascades through all themes, platforms, and 200+ components automatically.", VIOLET),
    ("📖", "Living Docs",            "Generates and refreshes Storybook stories and Figma annotations from actual usage.", YELLOW),
    ("🧩", "Variant Generator",      "Feed a base component + brand tokens. Agent outputs the full variant set, including edge cases.", CORAL),
    ("💬", "Natural Language Query", "Ask: 'Which buttons don't use primary-action token?' Get an exhaustive answer instantly.", GREEN),
    ("📋", "Changelog Writer",       "Reads every component commit and writes a human-friendly changelog entry, correctly categorised.", VIOLET),
    ("🔗", "Code Connect Sync",      "Keeps Figma component ↔ code mappings up to date as the codebase evolves.", PURPLE),
]

cols, rows = 4, 2
cw2, ch2 = Inches(3.0), Inches(1.55)
pad_x, pad_y = Inches(0.7), Inches(2.3)
gap_x, gap_y = Inches(0.18), Inches(0.18)

for i, (icon, label, desc, color) in enumerate(ideas):
    col, row = i % cols, i // cols
    x = pad_x + col * (cw2 + gap_x)
    y = pad_y + row * (ch2 + gap_y)
    card_rect(s3, x, y, cw2, ch2, fill=CARD)
    top_bar(s3, x, y, cw2, color, 4)
    txt(s3, icon + "  " + label, x+Pt(10), y+Pt(12), cw2-Pt(20), Pt(24), 14, bold=True, color=WHITE)
    txt(s3, desc, x+Pt(10), y+Pt(38), cw2-Pt(20), ch2-Pt(48), 11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 · HOW TO BUILD
# ═══════════════════════════════════════════════════════════════════════════
s4 = add_slide()
bg_rect(s4)
bar(s4, 0, 0, H, YELLOW)

eyebrow(s4, "How to Build", Inches(0.7), Inches(0.55), YELLOW)
txt(s4, "Five steps from zero to agentic",
    Inches(0.7), Inches(1.1), Inches(10), Inches(0.9),
    40, bold=True, color=WHITE)

steps = [
    ("1", "Audit your token layer",         "Agents need structured, semantic tokens to read intent. If your tokens are a mess, clean them first. This is the foundation everything else builds on.", YELLOW),
    ("2", "Connect Figma to code",           "Set up Code Connect so agents can see the canonical relationship between Figma components and their code implementations. Without this, agents are blind.", GREEN),
    ("3", "Deploy a drift-detection agent",  "Your first agent should have one verifiable job: compare Figma specs to code on every PR and comment when they diverge. Prove value fast.", VIOLET),
    ("4", "Add human-in-the-loop gates",     "Agents propose — humans approve. Build Slack/GitHub approval flows before letting agents auto-merge. Trust is earned incrementally.", CORAL),
    ("5", "Expand from proof of value",      "Once drift detection is trusted, layer in accessibility scanning, then token propagation, then docs generation. Each agent earns its place.", PURPLE),
]

sw, sh = Inches(11.93), Inches(0.82)
sx = Inches(0.7)
sy_start = Inches(2.15)
gap_s = Inches(0.1)

for i, (num, title, desc, color) in enumerate(steps):
    sy = sy_start + i * (sh + gap_s)
    fill = CARD if i % 2 == 0 else RGBColor(0x17,0x17,0x1E)
    card_rect(s4, sx, sy, sw, sh, fill=fill)
    top_bar(s4, sx, sy, Pt(5), color, int(sh / Pt(1)))  # left bar
    rect(s4, sx, sy, Pt(5), sh, color)                   # left accent
    # number badge circle (oval shape = 9)
    badge = s4.shapes.add_shape(9, sx+Pt(10), sy+Pt(12), Pt(36), Pt(36))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    txt(s4, num, sx+Pt(10), sy+Pt(14), Pt(36), Pt(32), 16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s4, title, sx+Pt(55), sy+Pt(10), Inches(2.8), Pt(30), 15, bold=True, color=WHITE)
    txt(s4, desc,  sx+Pt(55)+Inches(2.9), sy+Pt(10), Inches(8.5), sh-Pt(20), 12, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 · KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════════════════════
s5 = add_slide()
bg_rect(s5)
bar(s5, 0, 0, H, PURPLE)

eyebrow(s5, "Key Takeaways", Inches(0.7), Inches(0.55), VIOLET)
txt(s5, "What to take away from this",
    Inches(0.7), Inches(1.1), Inches(10), Inches(0.9),
    40, bold=True, color=WHITE)

takes = [
    ("Complexity has outpaced human bandwidth — agents aren't a trend, they're a necessity.", PURPLE),
    ("Clean token architecture is step zero. Garbage tokens mean confused agents.", VIOLET),
    ("Start with one narrow agent and one verifiable win. Scale from proof of value, not grand vision.", GREEN),
    ("Agents handle execution. You still own intent, quality bars, and the system's principles.", YELLOW),
    ("The question isn't if your design system becomes agentic — it's how prepared you'll be when it does.", CORAL),
]

tw, th2 = Inches(11.93), Inches(0.82)
tx2 = Inches(0.7)
ty_start = Inches(2.15)

for i, (text, color) in enumerate(takes):
    ty = ty_start + i * (th2 + gap_s)
    fill = CARD if i % 2 == 0 else RGBColor(0x17,0x17,0x1E)
    card_rect(s5, tx2, ty, tw, th2, fill=fill)
    rect(s5, tx2, ty, Pt(5), th2, color)
    # dot badge
    dot = s5.shapes.add_shape(9, tx2+Pt(12), ty+Pt(14), Pt(32), Pt(32))
    dot.fill.solid(); dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    txt(s5, str(i+1), tx2+Pt(12), ty+Pt(16), Pt(32), Pt(28), 14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s5, text, tx2+Pt(54), ty+Pt(18), tw-Pt(70), th2-Pt(24), 14, color=LIGHT)

txt(s5, "Build systems that build themselves.  ⚡",
    Inches(0.7), Inches(7.1), Inches(8), Pt(26),
    15, bold=True, color=VIOLET)


# ── Save ────────────────────────────────────────────────────────────────────
out = "/home/user/portfolio/Agentic_Design_Systems.pptx"
prs.save(out)
print("Saved:", out)
